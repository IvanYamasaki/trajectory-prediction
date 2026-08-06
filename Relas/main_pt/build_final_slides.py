from copy import deepcopy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "Seq2Seq Drift - Ivan Yamasaki.pptx"
OUT = HERE / "main_pt_slides.pptx"
FIG = HERE / "figs"

PURPLE = RGBColor(0x8D, 0x0D, 0xE3)
INK = RGBColor(0x0E, 0x0D, 0x0B)
BODY = RGBColor(0x5A, 0x58, 0x52)
MUTED = RGBColor(0x9E, 0x9B, 0x92)
SAND = RGBColor(0xD7, 0xD4, 0xCC)
PALE = RGBColor(0xF7, 0xF5, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def remove_all_slides(prs):
    # Keep the theme and package assets, discard the template's instructional pages.
    ids = list(prs.slides._sldIdLst)
    for sld_id in ids:
        rel_id = sld_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(sld_id)


def rgb(fill, color):
    fill.solid()
    fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def textbox(slide, x, y, w, h, text, size=12, color=BODY, bold=False,
            font="Google Sans", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
            margin=0.0, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for pi, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
    return shape


def rect(slide, x, y, w, h, color, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    rgb(shape.fill, color)
    no_line(shape)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=PURPLE, width=1.2):
    shp = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def source_icons(template):
    # The three small visual assets (Google Slides template signature) are reused.
    s = template.slides[3]
    return [BytesIO(sh.image.blob) for sh in s.shapes if sh.shape_type == 13]


def footer(slide, number, icons):
    for idx, blob in enumerate(icons):
        blob.seek(0)
        slide.shapes.add_picture(blob, Inches(0.35 + 0.26 * idx), Inches(5.17), height=Inches(0.16))
    textbox(slide, 9.35, 5.08, 0.35, 0.2, f"{number:02d}", 8, MUTED,
            align=PP_ALIGN.RIGHT)


def add_base(prs, title, subtitle, number, icons):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    rect(slide, 0, 0, 10, 5.625, WHITE)
    textbox(slide, 0.55, 0.53, 8.8, 0.38, title, 18, INK, True)
    textbox(slide, 0.55, 0.93, 8.8, 0.24, subtitle, 10, PURPLE)
    line(slide, 0.55, 1.28, 9.45, 1.28, SAND, 0.7)
    footer(slide, number, icons)
    return slide


def bullet_card(slide, x, y, w, h, label, text, accent=PURPLE):
    rect(slide, x, y, w, h, PALE, True)
    rect(slide, x + 0.18, y + 0.22, 0.07, h - 0.44, accent, True)
    textbox(slide, x + 0.42, y + 0.22, w - 0.6, 0.28, label, 10, accent, True)
    textbox(slide, x + 0.42, y + 0.61, w - 0.6, h - 0.78, text, 10.5, BODY)


def figure(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def title_slide(prs, icons):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    rect(slide, 0, 0, 10, 5.625, WHITE)
    # Same three-element top signature used by the supplied cover.
    for idx, blob in enumerate(icons):
        blob.seek(0)
        slide.shapes.add_picture(blob, Inches(3.3 + 1.7 * idx), Inches(1.25), height=Inches(0.55))
    textbox(slide, 1.05, 2.35, 7.9, 0.8,
            "Data Drift em um modelo Seq2Seq\npara predição de trajetórias na RoboCup SSL",
            22, INK, True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.45, 3.45, 7.1, 0.35,
            "Ivan G. Yamasaki · Orientador: Prof. Dr. Marcos R. O. A. Máximo",
            10.5, BODY, align=PP_ALIGN.CENTER)
    textbox(slide, 1.45, 3.84, 7.1, 0.28,
            "Instituto Tecnológico de Aeronáutica · Julho de 2026",
            10.5, MUTED, align=PP_ALIGN.CENTER)
    rect(slide, 4.35, 4.45, 1.3, 0.055, PURPLE, True)


def add_protocol(slide):
    # The method slide has three clear stages, following the template's methodology pattern.
    stages = [
        ("1 · Base", "Logs oficiais 2019–2025\n6 jogos da Divisão A por\nmundial · 36 no total\n2020 não teve competição"),
        ("2 · Stream", "8 mil janelas/jogo\n288 mil trajetórias\nordem temporal preservada"),
        ("3 · Avaliação", "ADE em milímetros\nhorizontes 30→15 e 60→30\n2019 = referência"),
    ]
    for i, (lab, desc) in enumerate(stages):
        x = 0.65 + i * 3.05
        rect(slide, x, 1.85, 2.65, 2.35, WHITE, True)
        line(slide, x, 1.85, x + 2.65, 1.85, PURPLE, 2.4)
        textbox(slide, x + 0.25, 2.15, 2.15, 0.3, lab, 11, PURPLE, True)
        textbox(slide, x + 0.25, 2.72, 2.1, 1.1, desc, 11, BODY)
    textbox(slide, 0.75, 4.48, 8.6, 0.28,
            "O baseline usa os seis jogos de 2019: com apenas um jogo, a ADE varia de 4,53 a 6,20 mm.",
            9.5, MUTED, italic=True)


def add_model(slide):
    items = [
        ("Entrada", "[x, y, vx, vy, sin ψ, cos ψ]\nψ = orientação do robô\n30 passos observados"),
        ("Seq2Seq", "Encoder LSTM (128)\natenção aditiva (Bahdanau et al.)\ndecoder auto-regressivo"),
        ("Saída", "15 posições futuras\nADE como métrica\nbaseline: filtro de Kalman\nde velocidade constante"),
    ]
    for i, (head, body) in enumerate(items):
        x = 0.7 + i * 3.0
        rect(slide, x, 1.9, 2.55, 2.05, PALE, True)
        textbox(slide, x + 0.22, 2.2, 2.1, 0.28, head, 11, PURPLE, True)
        textbox(slide, x + 0.22, 2.72, 2.1, 0.85, body, 10.5, BODY)
        if i < 2:
            textbox(slide, x + 2.58, 2.77, 0.35, 0.35, "→", 18, PURPLE, True,
                    align=PP_ALIGN.CENTER)
    textbox(slide, 1.35, 4.45, 7.4, 0.35,
            "O modelo prevê velocidades incrementais; por isso é diretamente sensível a mudanças cinemáticas.",
            10, BODY, align=PP_ALIGN.CENTER)


def add_adaptation_protocol(slide):
    bullet_card(slide, 0.75, 1.7, 2.7, 2.65, "PROTOCOLO COMUM",
                "Regime antigo: 2019\nDados novos: 2022–2025\n500 trajetórias por jogo\n2021 excluído (torneio em simulação)")
    bullet_card(slide, 3.65, 1.7, 2.7, 2.65, "ATUALIZAÇÃO LEVE",
                "Encoder congelado\n2 camadas do decoder treináveis\ntaxa 10⁻⁵ · 5 épocas")
    bullet_card(slide, 6.55, 1.7, 2.7, 2.65, "CRITÉRIO",
                "cat_ratio = fração das trajetórias antigas que pioram\n\nMenor é melhor: mede esquecimento catastrófico")


def add_detection_workflow(slide):
    steps = [
        ("Sinal", "ADE por trajetória\nnormalizada localmente\nmediana + MAD (w=200)"),
        ("Calibração", "ADWIN em grade de 96\n2-fold por jogos\nseleção e teste separados"),
        ("Ação", "alarme primário\n→ retreinar\ncamadas auxiliares\n→ monitorar e confirmar"),
    ]
    for i, (head, desc) in enumerate(steps):
        x = 0.65 + i * 3.1
        rect(slide, x, 1.9, 2.7, 2.25, PALE, True)
        textbox(slide, x + 0.25, 2.18, 2.15, 0.3, f"{i+1:02d}  {head}", 11, PURPLE, True)
        textbox(slide, x + 0.25, 2.76, 2.12, 1.1, desc, 10.5, BODY)
        if i < 2:
            textbox(slide, x + 2.73, 2.82, 0.32, 0.3, "→", 17, PURPLE, True)
    textbox(slide, 1.1, 4.57, 7.9, 0.26,
            "Pergunta operacional: o sinal é confiável o bastante para disparar retreinamento?",
            10, MUTED, align=PP_ALIGN.CENTER, italic=True)


def stat_card(slide, x, y, w, h, value, caption, unit=""):
    # Big-number card: readable from the back of the room, unlike a dense plot.
    rect(slide, x, y, w, h, PALE, True)
    textbox(slide, x + 0.2, y + 0.24, w - 0.4, 0.55, value, 26, PURPLE, True,
            align=PP_ALIGN.CENTER)
    if unit:
        textbox(slide, x + 0.2, y + 0.84, w - 0.4, 0.2, unit, 9, MUTED, True,
                align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.2, y + 1.06, w - 0.4, h - 1.2, caption, 9.5, BODY,
            align=PP_ALIGN.CENTER)


def add_detection_results(slide):
    # Replaces the three-panel figure, which was unreadable when projected.
    stat_card(slide, 0.7, 1.55, 2.72, 1.75, "0,042 – 0,125",
              "falsos alarmes em 2019, nos folds A e B\n(teto de admissibilidade: 0,20)",
              "POR MIL TRAJETÓRIAS")
    stat_card(slide, 3.64, 1.55, 2.72, 1.75, "9 / 10",
              "ano-folds com pelo menos um alarme\nna temporada em que houve drift",
              "COBERTURA")
    stat_card(slide, 6.58, 1.55, 2.72, 1.75, "0,2 – 1,5",
              "até o primeiro alarme depois da\nvirada de cada temporada",
              "JOGOS DE LATÊNCIA")
    bullet_card(slide, 0.7, 3.5, 4.2, 1.35, "SEM NORMALIZAÇÃO NÃO HÁ DETECTOR",
                "Nenhuma das 96 configurações é admissível no sinal bruto, e os 21 "
                "suavizadores testados pioram o FAR.")
    bullet_card(slide, 5.1, 3.5, 4.2, 1.35, "SEM HOLDOUT, OUTRA ESCOLHA",
                "A configuração vencedora dentro da amostra mantém FAR ≈ 0 e cobre "
                "só 2–3 de 5 temporadas fora dela.")


def add_detection_table(slide):
    cols = [0.7, 2.1, 4.55, 6.4, 7.85]
    widths = [1.35, 2.35, 1.65, 1.35, 1.3]
    headers = ["Camada", "Sinal", "FAR / 1k\n(A / B)", "Cobertura\n(A / B)", "Ação"]
    for x, w, h in zip(cols, widths, headers):
        rect(slide, x, 1.7, w, 0.62, PURPLE, True)
        textbox(slide, x + 0.08, 1.83, w - 0.16, 0.35, h, 8.7, WHITE, True, align=PP_ALIGN.CENTER)
    rows = [
        ("Primária", "ADE seq2seq", "0,042 / 0,125", "5/5 · 4/5", "retreinar"),
        ("Vigilância", "accel_p95", "0,292 / 0,250", "5/5 · 5/5", "monitorar"),
        ("Confirmação", "AND(ADE, accel)", "0,000 / 0,000", "3/5 · 2/5", "confirmar"),
        ("Corroboração", "ADE do filtro de Kalman", "0,167 / 0,083", "5/5 · 5/5", "redundância"),
    ]
    for r, row in enumerate(rows):
        y = 2.42 + r * 0.53
        bg = PALE if r % 2 == 0 else WHITE
        for x, w, val in zip(cols, widths, row):
            rect(slide, x, y, w, 0.48, bg, True)
            textbox(slide, x + 0.06, y + 0.12, w - 0.12, 0.2, val, 8.7,
                    PURPLE if x == cols[0] else BODY,
                    bold=(x == cols[0]), align=PP_ALIGN.CENTER)
    textbox(slide, 0.75, 4.85, 8.7, 0.22,
            "Resultados fora da amostra. A camada de aceleração antecipa a de erro em 8 de 10 folds anuais.",
            9.2, MUTED, italic=True, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation(TEMPLATE)
    icons = source_icons(prs)
    remove_all_slides(prs)

    title_slide(prs, icons)

    s = add_base(prs, "Introdução e motivação", "Predição precisa hoje; comportamento da liga muda amanhã", 2, icons)
    bullet_card(s, 0.7, 1.7, 2.7, 2.6, "CONTEXTO", "A predição de trajetórias apoia decisões em tempo real na Small Size League.")
    bullet_card(s, 3.65, 1.7, 2.7, 2.6, "LACUNA", "Um modelo treinado em uma temporada pode degradar silenciosamente em ligas futuras.")
    bullet_card(s, 6.6, 1.7, 2.7, 2.6, "OBJETIVO", "Medir o drift, explicar sua origem e definir como adaptar e quando retreinar.")

    s = add_base(prs, "Perguntas do estudo", "Uma avaliação temporal do preditor seq2seq usado como baseline", 3, icons)
    questions = [
        ("Quanto?", "Como a ADE muda entre 2019 e 2025?"),
        ("Por quê?", "Que parte vem de covariate shift e quais variáveis o explicam?"),
        ("O que fazer?", "Como adaptar sem esquecer e como sinalizar o momento de retreinar?"),
    ]
    for i, (h, b) in enumerate(questions):
        bullet_card(s, 0.75 + i * 2.95, 1.85, 2.65, 2.25, h, b)

    s = add_base(prs, "Dados e protocolo", "Logs oficiais dos mundiais da RoboCup SSL", 4, icons)
    add_protocol(s)

    s = add_base(prs, "Modelo e baseline", "Comparação limpa entre um modelo aprendido e um comparador cinemático fixo", 5, icons)
    add_model(s)

    s = add_base(prs, "Resultado 1 — o drift é real", "A perda de precisão cresce ao se afastar da temporada de treino", 6, icons)
    figure(s, FIG / "_fig_degradacao_pt.png", 0.75, 1.62, 8.5, 2.85)
    textbox(s, 0.85, 4.68, 8.35, 0.28, "São dois efeitos separados: o salto de 2021 (torneio em simulação) e a tendência gradual de 2022 a 2025, ambos medidos contra 2019.", 9.5, MUTED, italic=True, align=PP_ALIGN.CENTER)

    s = add_base(prs, "Resultado 2 — a aceleração lidera a mudança", "A cauda cinemática explica onde o novo regime entra", 7, icons)
    figure(s, FIG / "_fig_drift_cinematico_pt.png", 0.72, 1.62, 8.6, 2.78)
    textbox(s, 0.75, 4.67, 8.5, 0.32, "KS e Wasserstein se afastam da base 2019; accel_p99 se associa ao aumento da ADE por jogo.", 9.5, MUTED, italic=True, align=PP_ALIGN.CENTER)

    s = add_base(prs, "Resultado 3 — predominância de covariate shift", "Reponderar exemplos de 2019 recupera a maior parte do excesso de erro", 8, icons)
    figure(s, FIG / "_fig_iw_pt.png", 1.15, 1.7, 7.7, 3.05)
    textbox(s, 0.8, 4.85, 8.4, 0.22, "LSIF atribui 59–72% do excesso a covariate shift; classificador independente confirma mudança em todas as temporadas.", 9.3, MUTED, italic=True, align=PP_ALIGN.CENTER)

    s = add_base(prs, "Adaptação ao drift", "O objetivo é recuperar o regime novo sem perder desempenho no regime antigo", 9, icons)
    add_adaptation_protocol(s)

    s = add_base(prs, "Resultado 4 — replay ponderado é a melhor resposta", "Restringir a adaptação ao recorte mais extremo da aceleração piora o esquecimento", 10, icons)
    figure(s, FIG / "_fig_finetune_pt.png", 0.65, 1.62, 8.7, 2.65)
    textbox(s, 0.75, 4.62, 8.5, 0.40, "Fine-tuning leve reduz cat_ratio de 0,82 para ~0,46. Replay ponderado por importance weighting atinge o menor valor (0,417); EWC é inerte.", 9.4, MUTED, italic=True, align=PP_ALIGN.CENTER)

    s = add_base(prs, "Detecção online", "Do diagnóstico estatístico à decisão operacional de retreinar", 11, icons)
    add_detection_workflow(s)

    s = add_base(prs, "Arquitetura de detecção em camadas", "Baseada exclusivamente no protocolo e nos resultados do artigo", 12, icons)
    add_detection_table(s)

    s = add_base(prs, "Resultado 5 — detector primário confiável", "ADWIN sobre o erro normalizado alcança baixa taxa de alarmes fora da amostra", 13, icons)
    add_detection_results(s)
    textbox(s, 0.8, 4.93, 8.4, 0.2, "Todos os números são fora da amostra: a calibração usa metade dos jogos de cada ano e a outra metade é avaliada uma única vez.", 9.3, MUTED, italic=True, align=PP_ALIGN.CENTER)

    s = add_base(prs, "Conclusões", "O que o estudo entrega", 14, icons)
    bullet_card(s, 0.7, 1.65, 2.72, 2.8, "MEDIR", "Drift persistente após o treino. O baseline do filtro de Kalman piora ainda mais, porque os robôs se afastam do modelo de velocidade constante.")
    bullet_card(s, 3.64, 1.65, 2.72, 2.8, "EXPLICAR", "59–72% do excesso é covariate shift, com aceleração na cauda como vetor dominante.")
    bullet_card(s, 6.58, 1.65, 2.72, 2.8, "AGIR", "Replay ponderado reduz esquecimento; ADWIN normalizado fornece gatilho mensurável para retreinamento.")

    s = add_base(prs, "Limitações e próximos passos", "Onde as conclusões devem ser interpretadas com cuidado", 15, icons)
    bullet_card(s, 0.7, 1.65, 2.75, 2.65, "LIMITAÇÕES", "2021 mistura o gap entre simulação e realidade com o drift sazonal.\n\nSó a Divisão A e seis jogos por temporada.\n\nADE é uma métrica single-mode.")
    bullet_card(s, 3.62, 1.65, 2.75, 2.65, "CUIDADO OPERACIONAL", "Normalização robusta pressupõe drift gradual.\n\nHá autocorrelação entre trajetórias.\n\nA fração residual de concept drift não foi isolada por variável.")
    bullet_card(s, 6.54, 1.65, 2.75, 2.65, "PRÓXIMOS PASSOS", "Retreinar o encoder.\n\nEscalar o replay ponderado.\n\nAvaliar preditores multimodais e mais temporadas.")

    s = add_base(prs, "Referências essenciais", "Trabalhos que estruturam o modelo, o diagnóstico e a detecção", 16, icons)
    refs = [
        "Steuernagel, L.; Máximo, M. R. O. A. Trajectory prediction for SSL robots using seq2seq neural networks. RoboCup 2022, LNCS 13561, Springer, 2023.",
        "Gama, J. et al. A survey on concept drift adaptation. ACM Computing Surveys, 2014.",
        "Bifet, A.; Gavaldà, R. Learning from time-changing data with adaptive windowing (ADWIN). SDM, 2007.",
        "Kanamori, T.; Hido, S.; Sugiyama, M. A least-squares approach to direct importance estimation (LSIF). JMLR, 2009.",
    ]
    for i, ref in enumerate(refs):
        textbox(s, 0.85, 1.65 + i * 0.72, 8.35, 0.45, ref, 10, BODY)
        line(s, 0.85, 2.15 + i * 0.72, 9.15, 2.15 + i * 0.72, SAND, 0.45)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    rect(slide, 0, 0, 10, 5.625, WHITE)
    for idx, blob in enumerate(icons):
        blob.seek(0)
        slide.shapes.add_picture(blob, Inches(3.3 + 1.7 * idx), Inches(1.25), height=Inches(0.55))
    textbox(slide, 1.0, 2.65, 8.0, 0.55, "Obrigado!", 25, INK, True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.0, 3.4, 8.0, 0.30, "Perguntas e discussão", 12, PURPLE, align=PP_ALIGN.CENTER)
    rect(slide, 4.35, 4.25, 1.3, 0.055, PURPLE, True)

    prs.core_properties.title = "Data Drift em um modelo Seq2Seq para predição de trajetórias na RoboCup SSL"
    prs.core_properties.subject = "Apresentação de Iniciação Científica"
    prs.core_properties.author = "Ivan G. Yamasaki"
    prs.save(OUT)
    print(f"Generated {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
